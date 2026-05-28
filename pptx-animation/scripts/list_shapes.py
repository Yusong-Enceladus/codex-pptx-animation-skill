#!/usr/bin/env python3
"""List slide shapes for building PowerPoint animation specs.

The reported python_index is python-pptx's one-based order. Use it for
orientation only; prefer shape names in animation specs.
"""

from __future__ import annotations

import argparse
import csv
import json
import re
import sys
from pathlib import Path

from pptx import Presentation


EMU_PER_INCH = 914400


def clean_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def emu_to_inch(value: int) -> float:
    return round(value / EMU_PER_INCH, 3)


def iter_shape_rows(pptx_path: Path, slide_numbers: set[int] | None) -> list[dict[str, object]]:
    deck = Presentation(pptx_path)
    rows: list[dict[str, object]] = []
    for slide_no, slide in enumerate(deck.slides, start=1):
        if slide_numbers and slide_no not in slide_numbers:
            continue
        for python_index, shape in enumerate(slide.shapes, start=1):
            text = clean_text(getattr(shape, "text", "")) if getattr(shape, "has_text_frame", False) else ""
            rows.append(
                {
                    "slide": slide_no,
                    "python_index": python_index,
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "left_in": emu_to_inch(shape.left),
                    "top_in": emu_to_inch(shape.top),
                    "width_in": emu_to_inch(shape.width),
                    "height_in": emu_to_inch(shape.height),
                    "has_text": bool(getattr(shape, "has_text_frame", False)),
                    "text": text,
                }
            )
    return rows


def write_csv(rows: list[dict[str, object]], output: Path | None) -> None:
    fieldnames = [
        "slide",
        "python_index",
        "shape_id",
        "name",
        "type",
        "left_in",
        "top_in",
        "width_in",
        "height_in",
        "has_text",
        "text",
    ]
    if output:
        output.parent.mkdir(parents=True, exist_ok=True)
        with output.open("w", newline="", encoding="utf-8") as handle:
            writer = csv.DictWriter(handle, fieldnames=fieldnames)
            writer.writeheader()
            writer.writerows(rows)
    else:
        writer = csv.DictWriter(sys.stdout, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)


def write_json(rows: list[dict[str, object]], output: Path | None) -> None:
    text = json.dumps(rows, ensure_ascii=False, indent=2)
    if output:
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(text + "\n", encoding="utf-8")
    else:
        print(text)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--slides", nargs="*", type=int, help="Optional 1-based slide numbers to inspect")
    parser.add_argument("--output", type=Path, help="CSV or JSON output path. Defaults to CSV on stdout.")
    parser.add_argument("--format", choices=["csv", "json"], help="Output format. Defaults from --output suffix.")
    args = parser.parse_args()

    slide_numbers = set(args.slides) if args.slides else None
    output_format = args.format
    if output_format is None:
        output_format = "json" if args.output and args.output.suffix.lower() == ".json" else "csv"

    rows = iter_shape_rows(args.pptx, slide_numbers)
    if output_format == "json":
        write_json(rows, args.output)
    else:
        write_csv(rows, args.output)


if __name__ == "__main__":
    main()
